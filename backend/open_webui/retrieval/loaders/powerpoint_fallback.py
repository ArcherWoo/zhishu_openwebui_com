from langchain_core.documents import Document


class PowerPointFallbackLoader:
    def __init__(self, file_path: str, filename: str):
        self.file_path = file_path
        self.filename = filename

    def load(self) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(self.file_path)
        sections = []

        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if getattr(shape, 'has_text_frame', False):
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)

            if texts:
                sections.append(f'# Slide {idx}\n\n' + '\n\n'.join(texts))

        content = '\n\n'.join(sections).strip()
        if not content:
            content = f'# {self.filename}\n\n(No extractable text found)'

        return [
            Document(
                page_content=content,
                metadata={
                    'source': self.filename,
                    'content_format': 'markdown',
                    'parsed_by': 'python-pptx-fallback',
                    'original_extension': self.filename.split('.')[-1].lower(),
                },
            )
        ]
